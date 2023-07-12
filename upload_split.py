import os
import collections
import collections.abc
from pptx import Presentation
from splitter_funcs import file_split_appender
import PyPDF2
import pytesseract
from PIL import Image
import pdf2image
import streamlit as st

def extract_ppts(ppts):
    full_text_dict = {}
    for ppt in ppts:
        full_text_dict[ppt] = ''
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            n = 0
            if full_text_dict[ppt] != '':
                full_text_dict[ppt] = full_text_dict[ppt] + '\n\n'
            # Go through each shape in the slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Print all text in the shape
                    for paragraph in shape.text_frame.paragraphs:
                        full_text_dict[ppt] = full_text_dict[ppt] + ' '
                        for run in paragraph.runs:
                            n += 1
                            full_text_dict[ppt] = full_text_dict[ppt] + ' ' + run.text + ('\n' if n==1 else '')

            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                # Extract text from the notes
                if notes_slide.notes_text_frame is not None:
                    for paragraph in notes_slide.notes_text_frame.paragraphs:
                        for run in paragraph.runs:
                            full_text_dict[ppt] = full_text_dict[ppt] + ' ' + run.text
    return full_text_dict


def image_pdf_convertor(out_list, pdf, output_placeholder, img_num_text):
    image_path = "./tmp/images"
    os.makedirs(image_path, exist_ok=True)
    if output_placeholder:
        output_placeholder.write(f'starting image conversion.. {img_num_text}')
    pages = pdf2image.convert_from_path(pdf, 500)

    image_counter = 1
    for page in pages:
        filename = os.path.join(image_path, "page_" + str(image_counter) + ".jpg")
        page.save(filename, 'JPEG')
        image_counter = image_counter + 1

    for i in range(1, image_counter):
        filename = os.path.join(image_path, "page_" + str(i) + ".jpg")
        text = str(pytesseract.image_to_string(Image.open(filename)))
        out_list.append(text)
        os.remove(filename)
    return out_list

# for ECON 6341
def get_split_pdf_ppt(doc_dir=None, image_dir=None, metadata_divisor="Chunk"):
    output_placeholder = st.sidebar.empty()
    full_text_dict, pdfs = {}, []
    if doc_dir:
        output_placeholder.text('Compiling docs..')
        pdfs = [doc_dir + '\\' + f for f in os.listdir(doc_dir) if f.endswith('.pdf')]
        ppts = [doc_dir + '\\' + f for f in os.listdir(doc_dir) if not f.endswith('.pdf')]
        
        full_text_dict = extract_ppts(ppts)
    
    if image_dir:
        output_placeholder.text('Beginning image compiling..')
        pdf_images = [image_dir + '\\' + f for f in os.listdir(image_dir) if f.endswith('.pdf')]
        out_list = []
        for index, pdf in enumerate(pdf_images):
            out_list = image_pdf_convertor(out_list, pdf, output_placeholder, f'{index+1} of {len(pdf_images)}')
        full_text_dict[image_dir.rpartition('//')[2]] = out_list
    else:
        return None, None

    output_placeholder.text('Beginning text splitting..')

    split_text, metadata = file_split_appender(full_text_dict = full_text_dict, pdf_name_list = pdfs, chunk_size=1500, overlap=200, metadata_divisor=metadata_divisor)

    output_placeholder.write(f'Splitting text finished, number of splits is {len(split_text)}')

    return split_text, metadata

